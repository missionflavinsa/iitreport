from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# -------------------------------
# Create Presentation (16:9)
# -------------------------------
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# -------------------------------
# Short & Long Answer Questions
# (NO MCQs)
# -------------------------------
questions = [
    "7. State the two laws of reflection of light.",
    "8. What is lateral inversion? Give one example related to a plane mirror.",
    "9. Differentiate between regular reflection and diffused reflection.",
    "10. What is a kaleidoscope? State one of its applications.",
    "11. Which cells in the human eye are sensitive to bright light and colours, "
    "and which are sensitive to dim light?",

    "12. Draw a labelled diagram of the human eye and explain the functions of "
    "its main parts.",
    "13. What is persistence of vision? Explain with an example.",
    "14. Define myopia and hypermetropia. How are they corrected?",
    "21. Define amplitude, frequency and time period.",
    "22. How is sound produced in human beings?",

    "23. Differentiate between music and noise.",
    "24. Why is thunder heard after lightning?",
    "25. Distinguish between ultrasonic and infrasonic sounds.",
    "26. Describe an experiment to show that sound cannot travel through vacuum.",
    "27. Explain the working of the human ear."
]

# -------------------------------
# Create Slides (5 Q per slide)
# -------------------------------
for i in range(0, len(questions), 5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

    # Background: Black
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Text box with fixed padding (no overflow)
    textbox = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.75),
        Inches(11.8), Inches(6.0)
    )

    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.clear()

    for q in questions[i:i + 5]:
        p = text_frame.add_paragraph()
        p.text = q
        p.font.name = "Times New Roman"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

# -------------------------------
# Save File
# -------------------------------
prs.save("Light_and_Sound_Short_Answer_Question_Bank.pptx")

print("PPT generated successfully!")